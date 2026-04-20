from pathlib import Path

from PIL import Image, ImageDraw, ImageFilter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "本科毕业论文预答辩_华伯阳.pptx"
ASSET_DIR = ROOT / "ppt_assets"
SCREENSHOT = ROOT / "image" / "论文初稿_AI" / "1776267744781.png"


SLIDE_W = 13.333333
SLIDE_H = 7.5

COLORS = {
    "bg": "FFFFFF",
    "panel": "F0F3F8",
    "panel2": "E8EDF5",
    "line": "C5CDD8",
    "text": "1A2A4A",
    "muted": "5A6A82",
    "green": "1A3A6B",   # 深海军蓝（主色）
    "blue": "2255A4",    # 中蓝（次色）
    "amber": "B8922A",   # 金色（强调）
    "red": "C0392B",     # 红色（警示）
    "white": "FFFFFF",
    "road": "D8E0EC",    # 浅蓝灰（装饰）
}


def rgb(hex_value):
    hex_value = hex_value.strip("#")
    return RGBColor(int(hex_value[0:2], 16), int(hex_value[2:4], 16), int(hex_value[4:6], 16))


def make_assets():
    ASSET_DIR.mkdir(exist_ok=True)
    if SCREENSHOT.exists():
        img = Image.open(SCREENSHOT).convert("RGB")
        w, h = img.size
        crop = img.crop((int(w * 0.03), int(h * 0.10), int(w * 0.97), int(h * 0.93)))
        crop.thumbnail((2200, 1250), Image.LANCZOS)
        crop.save(ASSET_DIR / "system_screenshot.jpg", quality=92)

        focus = img.crop((int(w * 0.15), int(h * 0.35), int(w * 0.84), int(h * 0.82)))
        focus.thumbnail((1800, 950), Image.LANCZOS)
        focus.save(ASSET_DIR / "lane_focus.jpg", quality=92)

    # 学术风格装饰：浅蓝灰细线网格
    route = Image.new("RGBA", (1600, 900), (0, 0, 0, 0))
    draw = ImageDraw.Draw(route)
    for i in range(12):
        alpha = max(12, 55 - i * 4)
        color = (26, 58, 107, alpha)
        x0, y0 = 400 + i * 90, 800 - i * 55
        x1, y1 = 1500 - i * 20, 100 + i * 30
        draw.line([(x0, y0), (950, 450), (x1, y1)], fill=color, width=max(1, 3 - i // 5))
    route = route.filter(ImageFilter.GaussianBlur(0.5))
    route.save(ASSET_DIR / "route_lines.png")


def set_fill(shape, color, transparency=0):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)
    if transparency:
        fill.transparency = transparency


def set_line(shape, color, width=1, transparency=0):
    line = shape.line
    line.color.rgb = rgb(color)
    line.width = Pt(width)
    if transparency:
        line.transparency = transparency


def add_textbox(slide, text, x, y, w, h, font_size=22, color="text", bold=False,
                align=PP_ALIGN.LEFT, font="Microsoft YaHei", line_spacing=1.05):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = rgb(COLORS[color])
    return box


def add_title(slide, title, kicker=None):
    if kicker:
        add_textbox(slide, kicker, 0.72, 0.42, 5.0, 0.22, 8, "muted", bold=True)
    add_textbox(slide, title, 0.70, 0.60, 9.5, 0.58, 26, "green", bold=True)
    # 深蓝色下划线
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.70), Inches(1.26), Inches(1.6), Inches(0.04))
    set_fill(line, COLORS["green"])
    set_line(line, COLORS["green"], 0)
    # 金色短线
    gold = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(2.35), Inches(1.26), Inches(0.45), Inches(0.04))
    set_fill(gold, COLORS["amber"])
    set_line(gold, COLORS["amber"], 0)


def add_page_num(slide, n):
    add_textbox(slide, f"— {n:02d} —", 12.0, 6.95, 0.9, 0.25, 9, "muted", bold=False, align=PP_ALIGN.RIGHT)


def add_footer(slide, text="简易视频监控下车辆危险状态识别与实时预警工具设计"):
    add_textbox(slide, text, 0.70, 7.00, 8.5, 0.22, 8, "muted")


def add_background(slide, prs, accent=True):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb(COLORS["bg"])
    # 顶部深蓝色横条
    top = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.22))
    set_fill(top, COLORS["green"])
    set_line(top, COLORS["green"], 0)
    # 顶部金色细线
    gold_line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0.22), prs.slide_width, Inches(0.04))
    set_fill(gold_line, COLORS["amber"])
    set_line(gold_line, COLORS["amber"], 0)
    # 底部细线
    bot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(7.28), prs.slide_width, Inches(0.022))
    set_fill(bot, COLORS["line"])
    set_line(bot, COLORS["line"], 0)
    if accent:
        if (ASSET_DIR / "route_lines.png").exists():
            slide.shapes.add_picture(str(ASSET_DIR / "route_lines.png"), Inches(7.8), Inches(0.55), width=Inches(5.2))
    add_footer(slide)


def rounded_panel(slide, x, y, w, h, fill="panel", line="line", radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    panel = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(panel, COLORS[fill])
    set_line(panel, COLORS[line], 1)
    return panel


def pill(slide, text, x, y, w, h, fill, color="white", fs=10):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shape, COLORS[fill])
    set_line(shape, COLORS[fill], 0)
    add_textbox(slide, text, x, y + 0.03, w, h - 0.03, fs, color, bold=True, align=PP_ALIGN.CENTER)
    return shape


def bullet_list(slide, items, x, y, w, h, fs=15, color="text", gap=0.38, accent="green"):
    for i, item in enumerate(items):
        yy = y + i * gap
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(yy + 0.09), Inches(0.08), Inches(0.08))
        set_fill(dot, COLORS[accent])
        set_line(dot, COLORS[accent], 0)
        add_textbox(slide, item, x + 0.18, yy, w - 0.18, 0.32, fs, color)


def stat_card(slide, value, label, x, y, w=2.05, h=1.05, accent="green"):
    rounded_panel(slide, x, y, w, h, "panel2")
    add_textbox(slide, value, x + 0.16, y + 0.20, w - 0.32, 0.36, 24, accent, bold=True)
    add_textbox(slide, label, x + 0.17, y + 0.65, w - 0.32, 0.25, 9, "muted")


def flow_node(slide, title, subtitle, x, y, w, accent="green"):
    rounded_panel(slide, x, y, w, 0.92, "panel2")
    add_textbox(slide, title, x + 0.18, y + 0.16, w - 0.36, 0.24, 12, "text", bold=True)
    add_textbox(slide, subtitle, x + 0.18, y + 0.48, w - 0.36, 0.24, 8.5, "muted")
    mark = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.055), Inches(0.92))
    set_fill(mark, COLORS[accent])
    set_line(mark, COLORS[accent], 0)


def arrow(slide, x1, y1, x2, y2, color="green"):
    shape = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    shape.line.color.rgb = rgb(COLORS[color])
    shape.line.width = Pt(2.0)
    shape.line.end_arrowhead = True
    return shape


# ── Slide 1: 封面 ──────────────────────────────────────────────────────────────
def add_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)
    # 左侧深蓝色竖条装饰
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.70), Inches(0.55), Inches(0.06), Inches(5.8))
    set_fill(bar, COLORS["green"])
    set_line(bar, COLORS["green"], 0)
    gold_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.76), Inches(0.55), Inches(0.025), Inches(5.8))
    set_fill(gold_bar, COLORS["amber"])
    set_line(gold_bar, COLORS["amber"], 0)

    add_textbox(slide, "本科毕业论文预答辩", 1.05, 0.72, 5.0, 0.32, 11, "muted", bold=False)
    add_textbox(slide, "简易视频监控下车辆危险状态识别\n与实时预警工具设计",
                1.05, 1.22, 8.0, 1.55, 32, "green", bold=True, line_spacing=0.92)
    add_textbox(slide, "基于 OpenCV 传统视觉算法与 Flask Web 工具链的低成本驾驶辅助原型",
                1.08, 3.05, 7.8, 0.32, 13, "muted")

    # 分隔线
    sep = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.05), Inches(3.58), Inches(5.5), Inches(0.018))
    set_fill(sep, COLORS["line"])
    set_line(sep, COLORS["line"], 0)

    add_textbox(slide, "汇报人：华伯阳  202283840002", 1.08, 3.82, 5.0, 0.28, 12, "text", bold=True)
    add_textbox(slide, "指导教师：郭萍    学院：计算机学院，软件学院", 1.08, 4.22, 5.2, 0.26, 11, "muted")
    add_textbox(slide, "南京信息工程大学", 1.08, 4.65, 4.0, 0.26, 11, "muted")
    pill(slide, "滨江楼 BS201 · 下午 1:30 · 8分钟", 1.05, 5.28, 3.5, 0.36, "amber", "white", 10)

    if (ASSET_DIR / "lane_focus.jpg").exists():
        pic = slide.shapes.add_picture(str(ASSET_DIR / "lane_focus.jpg"),
                                       Inches(7.20), Inches(1.20), width=Inches(5.45))
        set_line(pic, COLORS["line"], 1)
    add_page_num(slide, 1)


# ── Slide 2: 目录 ──────────────────────────────────────────────────────────────
def add_agenda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)
    add_title(slide, "汇报目录", "AGENDA")
    items = [
        ("01", "研究背景与意义", "问题痛点、低成本定位"),
        ("02", "系统总体设计", "四层架构、技术栈"),
        ("03", "核心算法实现", "车道线 / 车速 / 转向角"),
        ("04", "系统演示与测试结果", "界面截图、帧率、预警验证"),
        ("05", "遇到的困难", "精度局限、鲁棒性、标定适配"),
        ("06", "后续工作计划", "论文完成情况与一个月安排"),
    ]
    for i, (num, title, desc) in enumerate(items):
        col = i % 2
        row = i // 2
        x = 0.84 + col * 6.15
        y = 1.72 + row * 1.42
        pill(slide, num, x, y + 0.06, 0.48, 0.32, "green", "white", 10)
        add_textbox(slide, title, x + 0.65, y, 2.6, 0.3, 16, "text", bold=True)
        add_textbox(slide, desc, x + 0.65, y + 0.42, 4.8, 0.26, 11, "muted")
        sep = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                                     Inches(x), Inches(y + 0.88), Inches(5.5), Inches(0.012))
        set_fill(sep, COLORS["line"])
        set_line(sep, COLORS["line"], 0)
    add_page_num(slide, 2)


# ── Slide 3: 研究背景与意义 ────────────────────────────────────────────────────
def add_background_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)
    add_title(slide, "研究背景与意义", "WHY THIS TOPIC")
    # 痛点卡片
    rounded_panel(slide, 0.78, 1.62, 5.65, 2.55, "panel")
    add_textbox(slide, "问题痛点", 1.06, 1.88, 1.5, 0.3, 16, "text", bold=True)
    bullet_list(slide, [
        "交通事故频发，超速与急转弯是主要诱因",
        "现有 ADAS 依赖激光雷达/毫米波雷达，成本高昂",
        "难以在低成本场景（校园通勤车、私家车）普及",
    ], 1.08, 2.38, 4.95, 1.3, fs=12.5, gap=0.44)
    # 定位卡片
    rounded_panel(slide, 6.78, 1.62, 5.65, 2.55, "panel")
    add_textbox(slide, "本研究定位", 7.06, 1.88, 1.6, 0.3, 16, "text", bold=True)
    bullet_list(slide, [
        "仅用普通摄像头（百元级硬件）",
        "无需 GPU，普通 CPU 实时运行（≥30fps）",
        "面向低成本场景：校园通勤车、私家车",
    ], 7.08, 2.38, 4.95, 1.3, fs=12.5, gap=0.44, accent="amber")
    # 三个指标
    stat_card(slide, "百元级", "普通摄像头，无需雷达", 0.78, 4.55, 3.55, 1.05, "green")
    stat_card(slide, "≥30fps", "普通 CPU 实时处理", 4.65, 4.55, 3.55, 1.05, "blue")
    stat_card(slide, "全栈闭环", "算法→预警→Web 展示", 8.52, 4.55, 3.55, 1.05, "amber")
    add_page_num(slide, 3)


# ── Slide 4: 系统总体架构 ──────────────────────────────────────────────────────
def add_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)
    add_title(slide, "系统总体架构", "SYSTEM ARCHITECTURE")
    layers = [
        ("视频输入", "摄像头 / 本地视频文件", "blue"),
        ("视频采集线程", "CameraThread — 循环读帧、视频源切换", "green"),
        ("核心检测模块", "Detector：车道线 + LK光流测速 + 转向趋势角", "green"),
        ("预警 + 持久化", "阈值判断 → SQLite alerts.db", "amber"),
        ("Flask Web 服务", "MJPEG 视频流 + REST API → 浏览器界面", "red"),
    ]
    for i, (title, subtitle, accent) in enumerate(layers):
        y = 1.55 + i * 0.96
        rounded_panel(slide, 1.85, y, 7.8, 0.72, "panel2")
        mark = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                                      Inches(1.85), Inches(y), Inches(0.06), Inches(0.72))
        set_fill(mark, COLORS[accent])
        set_line(mark, COLORS[accent], 0)
        add_textbox(slide, title, 2.12, y + 0.16, 2.0, 0.26, 13, "text", bold=True)
        add_textbox(slide, subtitle, 4.28, y + 0.18, 5.0, 0.24, 11, "muted")
        if i < len(layers) - 1:
            arrow(slide, 5.75, y + 0.72, 5.75, y + 0.96, "green")
    # 技术栈 pill
    pill(slide, "Python 3.13", 1.85, 6.42, 1.55, 0.32, "panel2", "muted", 10)
    pill(slide, "OpenCV 4.9", 3.52, 6.42, 1.55, 0.32, "panel2", "muted", 10)
    pill(slide, "Flask", 5.19, 6.42, 1.0, 0.32, "panel2", "muted", 10)
    pill(slide, "SQLite", 6.31, 6.42, 1.0, 0.32, "panel2", "muted", 10)
    add_page_num(slide, 4)


# ── Slide 5: 核心算法一——车道线检测 ───────────────────────────────────────────
def add_algo_lane(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)
    add_title(slide, "核心算法一：车道线检测", "LANE DETECTION")
    steps = [
        ("① 颜色+边缘", "HLS阈值+Sobel→二值图"),
        ("② 透视变换", "梯形ROI→鸟瞰图"),
        ("③ 滑窗搜索", "每帧从中间重搜\n避免累积漂移"),
        ("④ 多项式拟合", "二阶多项式曲线"),
        ("⑤ 曲率验证", "曲率过大则丢弃"),
    ]
    for i, (title, desc) in enumerate(steps):
        x = 0.72 + i * 2.38
        rounded_panel(slide, x, 1.52, 2.05, 1.55, "panel2")
        pill(slide, f"0{i+1}", x + 0.72, 1.68, 0.55, 0.28, "green", "white", 9.5)
        add_textbox(slide, title, x + 0.14, 2.08, 1.75, 0.26, 11.5, "green", bold=True)
        add_textbox(slide, desc, x + 0.14, 2.42, 1.75, 0.48, 9.5, "muted", line_spacing=1.1)
        if i < len(steps) - 1:
            arrow(slide, x + 2.05, 2.28, x + 2.38, 2.28, "green")
    # 示意图：透视变换 + 滑窗搜索
    warp_img = ASSET_DIR / "diagram_warp.png"
    win_img  = ASSET_DIR / "diagram_sliding_window.png"
    if warp_img.exists():
        pic = slide.shapes.add_picture(str(warp_img), Inches(0.72), Inches(3.28), width=Inches(6.0))
        set_line(pic, COLORS["line"], 1)
    if win_img.exists():
        pic2 = slide.shapes.add_picture(str(win_img), Inches(6.95), Inches(3.28), width=Inches(5.65))
        set_line(pic2, COLORS["line"], 1)
    add_page_num(slide, 5)


# ── Slide 6: 核心算法二——车速估算 ─────────────────────────────────────────────
def add_algo_speed(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)
    add_title(slide, "核心算法二：车速估算", "SPEED ESTIMATION")
    # 左侧流程
    rounded_panel(slide, 0.78, 1.62, 5.55, 4.55, "panel")
    add_textbox(slide, "LK 光流法流程", 1.06, 1.88, 2.5, 0.3, 16, "text", bold=True)
    flow_steps = [
        "Shi-Tomasi 角点检测（ROI：图像下方 42% 区域）",
        "Lucas-Kanade 稀疏光流追踪帧间位移 Δy",
        "速度换算：v = median(Δy) × fps / ppm × 3.6",
        "5帧滑动均值平滑，限制最大值 250 km/h",
    ]
    bullet_list(slide, flow_steps, 1.08, 2.42, 4.85, 1.6, fs=12.5, gap=0.52)
    # 公式图
    rounded_panel(slide, 1.08, 4.38, 4.85, 0.88, "panel2")
    f_speed = ASSET_DIR / "formula_speed.png"
    if f_speed.exists():
        slide.shapes.add_picture(str(f_speed), Inches(1.12), Inches(4.42), width=Inches(4.75))
    # 右侧说明
    rounded_panel(slide, 6.68, 1.62, 5.65, 2.15, "panel2")
    add_textbox(slide, "参数说明", 6.96, 1.88, 1.4, 0.28, 14, "text", bold=True)
    bullet_list(slide, [
        "fps：当前视频帧率（实测约 30fps）",
        "pixels_per_meter：透视标定后的像素/米比例",
        "ROI 限定下方区域，减少远景噪声干扰",
    ], 6.96, 2.38, 4.95, 1.0, fs=12, gap=0.44, accent="blue")
    rounded_panel(slide, 6.68, 4.05, 5.65, 2.12, "panel2")
    add_textbox(slide, "精度说明", 6.96, 4.30, 1.4, 0.28, 14, "amber", bold=True)
    add_textbox(slide, "单目视觉测速依赖 pixels_per_meter 标定精度，绝对速度误差较大；当前用于趋势判断与相对变化检测，提供在线标定工具支持手动校准。",
                6.96, 4.72, 4.95, 1.0, 11.5, "muted", line_spacing=1.15)
    # 光流截图（占位符 or 实际图）
    optical_flow_img = ASSET_DIR / "screenshot_optical_flow.jpg"
    if optical_flow_img.exists():
        pic = slide.shapes.add_picture(str(optical_flow_img), Inches(0.78), Inches(5.62), width=Inches(5.55))
        set_line(pic, COLORS["line"], 1)
    else:
        ph = rounded_panel(slide, 0.78, 5.62, 5.55, 1.55, "panel2")
        add_textbox(slide, "【截图占位】光流追踪可视化\n运行系统后截图，保存为 ppt_assets/screenshot_optical_flow.jpg",
                    1.05, 5.88, 5.0, 0.88, 10, "muted", align=PP_ALIGN.CENTER, line_spacing=1.2)
    add_page_num(slide, 6)


# ── Slide 7: 核心算法三——转向趋势角 ───────────────────────────────────────────
def add_algo_steering(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)
    add_title(slide, "核心算法三：转向趋势角", "STEERING ANGLE")
    # 原理说明
    rounded_panel(slide, 0.78, 1.62, 7.55, 2.85, "panel")
    add_textbox(slide, "消失点偏移法", 1.06, 1.88, 2.2, 0.3, 16, "text", bold=True)
    bullet_list(slide, [
        "鸟瞰图中取左右车道线顶部端点的中点 → 消失点坐标",
        "消失点相对图像水平中心的偏移量 → 转向趋势",
        "公式：angle = clip(offset × 45, −90°, 90°)",
        "8帧滑动均值平滑，抑制单帧抖动",
    ], 1.08, 2.42, 6.85, 1.6, fs=12.5, gap=0.44)
    # 公式图
    rounded_panel(slide, 1.08, 4.72, 6.85, 0.72, "panel2")
    f_angle = ASSET_DIR / "formula_angle.png"
    if f_angle.exists():
        slide.shapes.add_picture(str(f_angle), Inches(1.12), Inches(4.76), width=Inches(6.75))
    # 消失点示意图
    vp_img = ASSET_DIR / "diagram_vanishing_point.png"
    if vp_img.exists():
        pic = slide.shapes.add_picture(str(vp_img), Inches(0.78), Inches(5.62), width=Inches(7.85))
        set_line(pic, COLORS["line"], 1)
    # 右侧预警触发
    rounded_panel(slide, 8.68, 1.62, 3.75, 2.85, "panel2")
    add_textbox(slide, "预警触发条件", 8.96, 1.88, 2.2, 0.3, 14, "text", bold=True)
    pill(slide, "超速预警", 8.96, 2.38, 1.55, 0.32, "red", "white", 10)
    add_textbox(slide, "speed > 80 km/h（可配置）", 8.96, 2.82, 3.0, 0.26, 11, "muted")
    pill(slide, "急转弯预警", 8.96, 3.28, 1.55, 0.32, "amber", "white", 10)
    rounded_panel(slide, 8.68, 4.72, 3.75, 1.45, "panel2")
    add_textbox(slide, "冷却机制", 8.96, 4.96, 1.2, 0.26, 13, "text", bold=True)
    add_textbox(slide, "3秒冷却时间，避免重复记录\n预警写入 SQLite alerts.db",
                8.96, 5.32, 3.0, 0.55, 11, "muted", line_spacing=1.15)
    add_page_num(slide, 7)


# ── Slide 8: 系统界面演示 ──────────────────────────────────────────────────────
def add_demo(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs, accent=False)
    add_title(slide, "系统界面演示", "SYSTEM DEMO")

    # 4张截图：2×2布局（左侧大图 + 右侧3小图）
    screenshots = [
        ("screenshot_main.jpg",      "主界面：实时视频 + HUD 叠加",   0.72, 1.52, 5.8,  3.55),
        ("screenshot_dashboard.jpg", "数据仪表盘 + 预警历史",          6.72, 1.52, 5.85, 1.68),
        ("screenshot_calibrate.jpg", "透视标定面板",                   6.72, 3.38, 2.85, 1.68),
        ("screenshot_config.jpg",    "参数配置面板",                   9.72, 3.38, 2.85, 1.68),
    ]
    for fname, label, x, y, w, h in screenshots:
        img_path = ASSET_DIR / fname
        if img_path.exists():
            pic = slide.shapes.add_picture(str(img_path), Inches(x), Inches(y), width=Inches(w))
            set_line(pic, COLORS["line"], 1)
        else:
            rounded_panel(slide, x, y, w, h, "panel2")
            add_textbox(slide, f"【截图占位】\n{label}\n→ 保存为 ppt_assets/{fname}",
                        x + 0.12, y + h * 0.25, w - 0.24, h * 0.55,
                        9, "muted", align=PP_ALIGN.CENTER, line_spacing=1.2)
        add_textbox(slide, label, x, y + h + 0.02, w, 0.22, 8.5, "muted")

    pill(slide, "1280×720 离线处理 > 56 fps", 0.72, 6.88, 3.2, 0.30, "green", "white", 9)
    add_page_num(slide, 8)


# ── Slide 9: 测试结果 ──────────────────────────────────────────────────────────
def add_testing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)
    add_title(slide, "测试结果", "TEST RESULTS")
    # 环境说明
    rounded_panel(slide, 0.78, 1.62, 11.75, 0.62, "panel2")
    add_textbox(slide, "测试环境：macOS · Python 3.13.2 · OpenCV 4.9.0 · 测试视频：3个 1280×720（Lane / challenge / project）",
                1.05, 1.78, 11.0, 0.28, 11.5, "muted")
    # 表格
    headers = ["指标", "结果", "说明"]
    rows = [
        ("离线处理帧率", "> 56 fps ✅", "高于 30fps 实时要求"),
        ("车道线检测", "规则道路效果良好", "清晰标线场景稳定"),
        ("预警触发", "正确触发并记录 ✅", "速度/角度超阈值写入 SQLite"),
        ("速度估算", "趋势可用", "绝对精度受标定影响"),
    ]
    col_x = [0.86, 3.55, 6.05]
    col_w = [2.45, 2.25, 5.65]
    for i, h in enumerate(headers):
        pill(slide, h, col_x[i], 2.52, col_w[i], 0.34, "green", "white", 10)
    for r, row in enumerate(rows):
        y = 3.14 + r * 0.72
        rounded_panel(slide, 0.78, y - 0.05, 11.0, 0.52, "panel2")
        add_textbox(slide, row[0], col_x[0] + 0.06, y + 0.04, col_w[0] - 0.12, 0.22, 10.5, "text", bold=True)
        ok = "green" if "✅" in row[1] else "amber"
        add_textbox(slide, row[1], col_x[1] + 0.06, y + 0.04, col_w[1] - 0.12, 0.22, 10.5, ok, bold=True)
        add_textbox(slide, row[2], col_x[2] + 0.06, y + 0.04, col_w[2] - 0.12, 0.22, 10.5, "muted")
    rounded_panel(slide, 0.78, 6.05, 11.0, 0.62, "panel")
    add_textbox(slide, "结论", 1.06, 6.24, 0.6, 0.22, 12, "green", bold=True)
    add_textbox(slide, "系统已完成从算法处理到浏览器展示的完整流程，可作为低成本车辆危险状态识别与预警原型验证工具。",
                1.74, 6.23, 9.2, 0.24, 12, "text")
    add_page_num(slide, 9)


# ── Slide 10: 遇到的困难 ───────────────────────────────────────────────────────
def add_difficulties(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)
    add_title(slide, "遇到的困难与解决思路", "RISKS & SOLUTIONS")
    cards = [
        ("速度估算精度有限",
         "pixels_per_meter 依赖实际场景标定，难以精确获取，绝对速度误差较大。",
         "当前用于趋势判断；提供在线标定工具，支持手动校准并在论文中说明局限性。"),
        ("复杂光照鲁棒性不足",
         "强光/阴影导致 HLS 阈值化失效，车道线丢失，影响检测连续性。",
         "加入帧间稳定性检验，使用上一帧结果作为备用，减少画面抖动。"),
        ("透视变换参数适配",
         "不同摄像头安装位置、焦距不同，需重新标定透视变换参数。",
         "开发了 Web 端在线拖拽标定工具，支持实时鸟瞰图预览与参数保存。"),
    ]
    for i, (title, problem, solution) in enumerate(cards):
        x = 0.78 + i * 3.75
        rounded_panel(slide, x, 1.80, 3.32, 4.05, "panel2")
        pill(slide, f"0{i+1}", x + 0.25, 2.05, 0.5, 0.32, "amber", "white", 10)
        add_textbox(slide, title, x + 0.90, 2.04, 2.0, 0.3, 14, "text", bold=True)
        add_textbox(slide, "困难", x + 0.28, 2.75, 0.65, 0.22, 10, "red", bold=True)
        add_textbox(slide, problem, x + 0.28, 3.08, 2.65, 0.72, 11, "muted", line_spacing=1.15)
        add_textbox(slide, "解决", x + 0.28, 4.18, 0.65, 0.22, 10, "green", bold=True)
        add_textbox(slide, solution, x + 0.28, 4.52, 2.65, 0.88, 11, "text", line_spacing=1.15)
    add_page_num(slide, 10)


# ── Slide 11: 论文完成情况 ─────────────────────────────────────────────────────
def add_thesis_progress(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)
    add_title(slide, "论文完成情况", "THESIS STATUS")
    done = [
        "第1章  绪论",
        "第2章  相关理论与技术基础",
        "第3章  系统总体设计",
        "第4章  核心算法设计与实现",
        "第5章  系统开发与实现",
        "第6章  系统测试与结果分析",
        "第7章  结论与展望",
    ]
    rounded_panel(slide, 0.78, 1.62, 5.55, 4.55, "panel")
    add_textbox(slide, "已完成（全部7章初稿）", 1.06, 1.88, 3.0, 0.3, 15, "green", bold=True)
    for i, ch in enumerate(done):
        y = 2.42 + i * 0.52
        pill(slide, "✅", 1.08, y + 0.04, 0.42, 0.28, "green", "white", 9)
        add_textbox(slide, ch, 1.62, y, 3.8, 0.28, 12, "text")
    # 待完善
    rounded_panel(slide, 6.68, 1.62, 5.65, 4.55, "panel2")
    add_textbox(slide, "待完善", 6.96, 1.88, 1.2, 0.3, 15, "amber", bold=True)
    todo = [
        ("第6章", "补充量化测试指标（帧率、检测率等）"),
        ("全文", "格式规范化（NUIST 标准）"),
        ("参考文献", "整理完善，补充近年文献"),
    ]
    for i, (tag, desc) in enumerate(todo):
        y = 2.42 + i * 0.72
        pill(slide, tag, 6.96, y + 0.04, 0.88, 0.30, "amber", "white", 10)
    rounded_panel(slide, 6.68, 5.08, 5.65, 1.09, "panel")
    add_textbox(slide, "当前阶段判断", 6.96, 5.28, 2.0, 0.26, 13, "text", bold=True)
    add_textbox(slide, "代码已全部完成，7章初稿已形成，主要工作为完善与规范化。",
                6.96, 5.62, 5.0, 0.38, 11.5, "muted", line_spacing=1.1)
    add_page_num(slide, 11)


# ── Slide 12: 后续工作计划 ─────────────────────────────────────────────────────
def add_plan(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)
    add_title(slide, "后续工作计划", "NEXT 4 WEEKS")
    weeks = [
        ("4.19–4.25", "补充测试数据", "完善 Lane / challenge / project 视频测试，整理量化结果图表"),
        ("4.26–5.2",  "全文格式规范", "统一术语、补齐系统设计图、完善结果分析与局限性说明"),
        ("5.3–5.9",   "修改润色",     "全文修改润色，提交导师审阅，根据反馈调整"),
        ("5.10–5.16", "定稿答辩",     "完成查重修改、PPT演练、答辩问题准备和最终材料提交"),
    ]
    for i, (week, title, desc) in enumerate(weeks):
        y = 1.60 + i * 1.04
        pill(slide, week, 0.88, y + 0.12, 1.35, 0.34, "green" if i < 2 else "amber", "white", 10)
        add_textbox(slide, title, 2.72, y + 0.17, 1.5, 0.24, 13, "text", bold=True)
        add_textbox(slide, desc, 4.38, y + 0.17, 6.8, 0.24, 11.5, "muted")
    rounded_panel(slide, 0.88, 5.95, 10.75, 0.62, "panel")
    add_textbox(slide, "可行性", 1.15, 6.14, 0.85, 0.22, 12, "green", bold=True)
    add_textbox(slide, "代码已全部完成，论文7章初稿已完成，主要工作为完善和规范化，一个月内可高质量完成。",
                2.08, 6.13, 8.7, 0.24, 11.5, "text")
    add_page_num(slide, 12)


# ── Slide 13: 总结 ─────────────────────────────────────────────────────────────
def add_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)
    add_textbox(slide, "总结", 0.80, 0.86, 1.3, 0.45, 26, "green", bold=True)
    add_textbox(slide, "提出基于普通摄像头的低成本车辆危险状态识别方案，实现车道线检测、车速估算、转向趋势角三合一系统，完成从算法到 Web 全栈的完整工程实现，系统在普通 CPU 上实时运行帧率超过 56fps。",
                0.82, 1.62, 8.6, 0.88, 20, "text", bold=True, line_spacing=1.05)
    # 贡献卡片
    contribs = [
        ("低成本方案", "仅用普通摄像头，无需雷达或 GPU"),
        ("三合一系统", "车道线 + 车速 + 转向角一体化"),
        ("全栈实现", "算法→预警→Web 完整闭环"),
        ("> 56fps", "普通 CPU 满足实时性要求"),
    ]
    for i, (val, label) in enumerate(contribs):
        x = 0.82 + i * 2.98
        stat_card(slide, val, label, x, 3.05, 2.65, 1.05, "green" if i % 2 == 0 else "blue")
    rounded_panel(slide, 0.82, 4.55, 4.95, 1.28, "panel2")
    add_textbox(slide, "下一步关键词", 1.12, 4.80, 1.6, 0.25, 13, "muted")
    add_textbox(slide, "测试补充 · 论文规范 · 演示稳定", 1.12, 5.22, 3.8, 0.28, 18, "amber", bold=True)
    add_textbox(slide, "谢谢各位老师！请批评指正。", 0.84, 6.18, 5.5, 0.45, 24, "text", bold=True)
    add_page_num(slide, 13)


# ── Build ──────────────────────────────────────────────────────────────────────
def build():
    make_assets()
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    add_cover(prs)           # 1
    add_agenda(prs)          # 2
    add_background_slide(prs)  # 3
    add_architecture(prs)    # 4
    add_algo_lane(prs)       # 5
    add_algo_speed(prs)      # 6
    add_algo_steering(prs)   # 7
    add_demo(prs)            # 8
    add_testing(prs)         # 9
    add_difficulties(prs)    # 10
    add_thesis_progress(prs) # 11
    add_plan(prs)            # 12
    add_closing(prs)         # 13

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
